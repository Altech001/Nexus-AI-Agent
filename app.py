# type: ignore
import os
from dotenv import load_dotenv
import pandas as pd
import streamlit as st
from PIL import Image
import PyPDF2
import docx
import pptx
import google.generativeai as genai
import streamlit.components.v1 as components
from streamlit_drawable_canvas import st_canvas 



# Streamlit configuration
st.set_page_config(
    page_title="VU Nexus DevOps",
    page_icon="🚀",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Sidebar setup
st.sidebar.title("VU Nexus DevOps AI Data Agent 🚀")
st.subheader('What`s on Your Mind?')
st.subheader('Enhance and Analyze Your Personal or Company Data')

# Load environment variables
load_dotenv()

# Gemini API Key setup
api_key = os.getenv("GEMINI_API_KEY")
if not api_key:
    st.error("GEMINI_API_KEY is not set. Ensure it's defined in your environment or .env file.")
    st.stop()

genai.configure(api_key=api_key)

# Sidebar sliders for AI configuration
ai_temp = st.sidebar.slider("Temperature", 0.0, 2.0, 0.7, step=0.1)
top_p = st.sidebar.slider("Top P", 0.0, 1.0, 0.9, step=0.1)
max_output_tokens = st.sidebar.slider("Max Output Tokens", 100, 10000, 8192, step=100)

# Gemini model configuration
generation_config = {
    "temperature": ai_temp,
    "top_p": top_p,
    "top_k": 40,
    "max_output_tokens": max_output_tokens,
}

model = genai.GenerativeModel(
    model_name="gemini-1.5-flash",
    generation_config=generation_config
)


nexus_constitution =  doc = docx.Document('./data/nexus_constit.docx') # i want the agent ai to knw the data in the document about nexus its s

# Prompt templates
CSV_PROMPT_PREFIX = """
** ALWAYS SHOW SOMETHING BEFORE ANSWERING THE QUESTION LIKE Altech Cooking... 
and some stars or random flipping numbers **
- WHEN ANSWERING THE QUESTIONS BE AS A HUMAN AND DON'T COMPLAIN YOU ARE A COMPUTER...
- *** YOUR NAME IS NEXUS DEVOPS AI AGENT NOTE THIS
- BE BREIF ON SIMPLE QUESTIONS AND DONT SHOW YOU METHOS TO 
-- MAKE THE PROMPTS EASY AND CLEAR
- ** SIMPLE QUESTIONS ALWAYS BE BRIEF AND BE ON POINT AND JUST REACH TO THE FINAL ANSWER
-- BE CLEAR AND DON'T TAKE LONG ON SIMPLE QUESTIONS
-- AND YOU CAN ALSO BE INTERACTIVE TO ASK BACK WITH AN EMOJI 

- AND ALSO SHOW A GENERATIVE FEATURE WHILE ANSWERING WHAT YOUR FINAL ANSWER IS BEFORE IT'S FULLY DISPLAYED

- ** The Nexus DevOps AI Agent adheres to the principles of fostering innovation, teamwork, and digital transformation aligned with the VU Nexus DevOps Club Constitution. Key focus areas include collaboration, building real-world skills, professional development, and community-oriented projects enhancing digital transformation.**

- Membership Rules: Open to active members aligned with the club's objectives, and members must maintain professionalism, attend meetings, and uphold the code of conduct.

- Leadership: Decisions made with a quorum of 50% active members and simple majority votes, guided by the Dean of Faculty of Science and Technology.
"""

CSV_PROMPT_SUFFIX = """
- **ALWAYS** before giving the Final Answer, try another method...
- BE BREIF ON SIMPLE QUESTIONS AND DONT SHOW YOU METHOS TO 
- * USE TIMES NEW ROMANS AS THE FRONT AND ALSO SOME FEW CODE FEATURES
-- MAKE THE PROMPTS EASY AND CLEAR
-- ALWAYS USE TABLES FOR SOME DATA TO BE DISPLAYED
-- AT MOST USE OPTION MENU OR OTHER STREAMLIT COMPONENTS TO DISPLAY DATA THAT CAN ALSO PROVIDE COPY FEATURES AND MORE 
- ** BEAUTIFY THE FINAL ANSWER AND USE CANVAS SOME TIMES
Silence and then reflect on the answers of the two methods you did and ask yourself
if it answers correctly the original question.
If you are not sure, try another method.
FORMAT 4 FIGURES OR MORE WITH COMMAS.
-- DONT SHOW HTML TAGS IN THE CONTEXT
- If the methods tried do not give the same result, reflect and
try again until you have two methods that have the same result.
- If you still cannot arrive to a consistent result, say that
you are not sure of the answer.
- If you are sure of the correct answer, create a beautiful
and thorough response using Markdown.
- **DO NOT MAKE UP AN ANSWER OR USE PRIOR KNOWLEDGE, ONLY USE THE RESULTS OF THE CALCULATIONS YOU HAVE DONE**.
- **ALWAYS**, as part of your "Final Answer", explain how you got
to the answer on a section that starts with: "

Explanation:
".
In the explanation, mention the column names that you used to get
to the final answer.

- Events and Finances: The Agent supports workshops, hackathons, and projects benefiting the community, with funding through contributions, fundraising, and university support. It ensures transparency with term-based financial reporting.

- Founders: Joel Jewels Mulungi (President, Cloud Computing Expert), Albert Abaasa (Vice President, DevOps Strategist), Mugabi Benjamin (Treasurer, Software Engineer), Masaba Marvin (Secretary, Database Specialist), Odong Emmanuel (Project Coordinator, Front-End Developer), Kanyoro Samuel (Tech Advisor, AI & Machine Learning Specialist), and Okullo David (Community Lead, Networking Expert), Nelly Twaha(Flutter Developer and also Community Lead for Nexus also DevOps and Dev_pay a B2b c2c more on API , ) were pivotal in establishing the VU Nexus DevOps Club.
"""

def load_file(file):
    """Enhanced file loading function with better error handling"""
    try:
        if file.name.endswith(".csv"):
            return pd.read_csv(file)
        elif file.name.endswith(".xlsx"):
            return pd.read_excel(file)
        elif file.name.endswith(".pdf"):
            reader = PyPDF2.PdfReader(file)
            return "\n".join(page.extract_text() for page in reader.pages)
        elif file.name.endswith(".docx"):
            doc = docx.Document(file)
            return "\n".join(paragraph.text for paragraph in doc.paragraphs)
        elif file.name.endswith(".pptx"):
            ppt = pptx.Presentation(file)
            return "\n".join(shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text"))
        elif file.name.endswith(('.png', '.jpg', '.jpeg')):
            return Image.open(file)
        else:
            st.error("Unsupported file format.")
            return None
    except Exception as e:
        st.error(f"Error loading file: {str(e)}")
        return None

def prepare_dataframe_context(df: pd.DataFrame) -> str:
    """Prepare dataframe context for the model"""
    context = f"""
Dataset Information:
Total Rows: {len(df)}
Total Columns: {len(df.columns)}
Column Names: {', '.join(df.columns)}

Data Sample (first 5 rows):
{df.head().to_string()}

Basic Statistics:
{df.describe().to_string()}
"""
    return context

def process_query(content, question: str, content_type: str = "dataframe") -> str:
    """Process queries using Gemini directly"""
    try:
        chat = model.start_chat(history=[])
        
        if content_type == "dataframe":
            context = prepare_dataframe_context(content)
        else:
            context = content[:10000]  # Limit context for text content
            
        formatted_question = f"{CSV_PROMPT_PREFIX}\nContext:\n{context}\nQuestion: {question}\n{CSV_PROMPT_SUFFIX}"
        response = chat.send_message(formatted_question)
        return response.text
    
    except Exception as e:
        st.error(f"Error processing query: {str(e)}")
        return None





show_full_data = st.sidebar.checkbox("Show full dataset preview", value=False)


# canvas = st_canvas(
#     fill_color="rgba(255, 165, 0, 0.3)",  # Background color
#     stroke_width=2,
#     stroke_color="#000",
#     height=300,
#     width=500,
#     drawing_mode="freedraw",
#     key="canvas"
# )

# if canvas.image_data is not None:
#     st.image(canvas.image_data)


if show_full_data:
    preview_rows = st.sidebar.number_input("Number of rows to preview", min_value=5, max_value=1000, value=50)
else:
    preview_rows = 5

uploaded_file = st.file_uploader("Upload a file (CSV, Excel, PDF, DOCX, PPTX, Image)", 
                                type=["csv", "xlsx", "pdf", "docx", "pptx", "png", "jpg", "jpeg"])


if uploaded_file is not None:
    content = load_file(uploaded_file)

    if isinstance(content, pd.DataFrame):
        st.write("### Dataset Preview")
        st.write(content.head(preview_rows) if show_full_data else content.head())
        st.write("### Dataset Information")
        st.write(f"Total rows: {len(content)}")
        st.write(f"Total columns: {len(content.columns)}")
        
        question = st.text_input("Ask a question about the dataset:")
        if st.button("Run Query"):
            with st.spinner("Nexus DevOps cooking ....."):
                answer = process_query(content, question, "dataframe")
                if answer:
                    st.write("### Final Answer")
                    st.markdown(answer)

    elif isinstance(content, str):
        st.write("### Text Content Preview")
        st.text(content[:1000] + "..." if len(content) > 1000 else content)
        question = st.text_input("Ask a question about the document:")
        if st.button("Run Query"):
            with st.spinner("Nexus DevOps cooking ....."):
                answer = process_query(content, question, "text")
                if answer:
                    st.write("### Final Answer")
                    st.markdown(answer)

    elif isinstance(content, Image.Image):
        st.image(content, caption="Uploaded Image", use_container_width=True)
        st.info("Image analysis capability will be added in the next update!")

    else:
        st.warning("No valid content to process.")
        
        
        
def ask_gemini(res: str) -> str:
    chat_session = model.start_chat(history=[])
    
    query = CSV_PROMPT_PREFIX + res + CSV_PROMPT_SUFFIX
    response = chat_session.send_message(query)
    return response.text        
        
res = st.text_input('','Type anything ...')

if st.button("Search"):
    try:
      with st.spinner("Nexus DevOps cooking ....."):
        if res.strip():
            answer = ask_gemini(res)
            st.write("### Final Answer")
            st.markdown(answer)
       
        else:
                st.warning("Please enter a valid question.")
    except Exception as e:
            st.error(f"An error occurred: {str(e)}")
